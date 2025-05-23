import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract
import os

def parse_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

def parse_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_image(file_path):
    image = Image.open(file_path)
    return pytesseract.image_to_string(image)

def parse_file(file_path, ext):
    try:
        if ext == ".pdf":
            return parse_pdf(file_path)
        elif ext == ".docx":
            return parse_docx(file_path)
        elif ext == ".pptx":
            return parse_pptx(file_path)
        elif ext in [".jpeg", ".jpg", ".png"]:
            return parse_image(file_path)
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Error parsing file: {str(e)}"
